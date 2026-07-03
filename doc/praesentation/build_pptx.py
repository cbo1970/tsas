#!/usr/bin/env python3
"""Generates the TSaS project presentation (architecture, deployment, AI insights)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----------------------------------------------------------------
INK      = RGBColor(0x1A, 0x23, 0x2E)   # near-black text
MUTED    = RGBColor(0x5B, 0x66, 0x70)   # secondary text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
BG       = RGBColor(0xF4, 0xF6, 0xF8)
ACCENT   = RGBColor(0x1B, 0x9E, 0x57)   # tennis green
ACCENT_D = RGBColor(0x14, 0x6B, 0x3C)
BLUE     = RGBColor(0x15, 0x65, 0xC0)
BLUE_BG  = RGBColor(0xBB, 0xDE, 0xFB)
GREEN_BG = RGBColor(0xC8, 0xE6, 0xC9)
YELL_BG  = RGBColor(0xFF, 0xF3, 0xC4)
ORNG_BG  = RGBColor(0xFF, 0xE0, 0xB2)
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
LINE     = RGBColor(0xD3, 0xDA, 0xE0)
WARN     = RGBColor(0xC6, 0x28, 0x28)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _set_font(run, size, color, bold, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (txt,size,color,bold) tuples."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            _set_font(r, size, color, bold)
    return tb


def card(s, l, t, w, h, fill=CARD, line=LINE, line_w=0.75, radius=0.06, shadow=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if shadow:
        el = sh._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {}); el.append(ef)
        sd = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad': '50000', 'dist': '25000', 'dir': '5400000', 'rotWithShape': '0'})
        ef.append(sd)
        clr = sd.makeelement(qn('a:srgbClr'), {'val': '1A232E'}); sd.append(clr)
        alpha = clr.makeelement(qn('a:alpha'), {'val': '18000'}); clr.append(alpha)
    return sh


def header(s, kicker, title):
    # top accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, Inches(0.6), Inches(0.42), Inches(12.1), Inches(0.35),
         [[(kicker.upper(), 12, ACCENT_D, True)]])
    text(s, Inches(0.6), Inches(0.72), Inches(12.1), Inches(0.7),
         [[(title, 28, INK, True)]])
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.42),
                            Inches(1.1), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    ln.shadow.inherit = False


def chip(s, l, t, w, h, label, fill, txt_color=WHITE, size=11, bold=True, line=None):
    c = card(s, l, t, w, h, fill=fill, line=line, radius=0.18)
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(label):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(0); p.space_before = Pt(0)
        for (txt, sz, col, bd) in para:
            r = p.add_run(); r.text = txt; _set_font(r, sz, col, bd)
    return c


def bullets(s, l, t, w, h, items, size=14, gap=7, color=INK, marker="—",
            marker_color=ACCENT):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_top = 0; tf.margin_right = 0; tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = 1.05
        rm = p.add_run(); rm.text = marker + "  "; _set_font(rm, size, marker_color, True)
        if isinstance(it, tuple):
            head, rest = it
            r1 = p.add_run(); r1.text = head; _set_font(r1, size, color, True)
            if rest:
                r2 = p.add_run(); r2.text = rest; _set_font(r2, size, MUTED, False)
        else:
            r = p.add_run(); r.text = it; _set_font(r, size, color, False)
    return tb


def connector(s, x1, y1, x2, y2, color=MUTED, width=1.5, dash=False):
    cn = s.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    cn.shadow.inherit = False
    if dash:
        ln = cn.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'}); ln.append(d)
    return cn


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s = slide()
panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
panel.fill.solid(); panel.fill.fore_color.rgb = INK; panel.line.fill.background()
panel.shadow.inherit = False
# accent side band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), EMU_H)
band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
band.shadow.inherit = False
text(s, Inches(0.95), Inches(1.7), Inches(11), Inches(0.5),
     [[("TENNIS SCORE AND STATISTIC", 15, RGBColor(0x8B, 0xE0, 0xB0), True)]])
text(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(1.6),
     [[("TSaS – Architektur,", 46, WHITE, True)],
      [("Deployment & KI-Erkenntnisse", 46, WHITE, True)]], space_after=2)
text(s, Inches(0.95), Inches(4.4), Inches(11), Inches(0.6),
     [[("Web-Applikation zur Tennismatch-Dokumentation und Statistik", 17,
        RGBColor(0xB9, 0xC2, 0xCC), False)]])
# meta chips
for i, (lab) in enumerate(["Spring Boot 4 · Java 25", "Angular · Keycloak",
                           "PostgreSQL · Docker", "Spring AI · OpenAI"]):
    chip(s, Inches(0.95 + i * 2.75), Inches(5.5), Inches(2.55), Inches(0.55),
         [[(lab, 11.5, WHITE, True)]], fill=RGBColor(0x2A, 0x37, 0x45),
         line=RGBColor(0x3D, 0x4C, 0x5B))
text(s, Inches(0.95), Inches(6.6), Inches(11), Inches(0.4),
     [[("Clean Architecture · Modularer Monolith · arc42", 12,
        RGBColor(0x8A, 0x96, 0xA2), False)]])

# =============================================================================
# SLIDE 2 — Überblick
# =============================================================================
s = slide()
header(s, "Projekt-Überblick", "Was ist TSaS?")
text(s, Inches(0.6), Inches(1.7), Inches(7.0), Inches(1.6),
     [[("Coaches und Eltern dokumentieren Tennismatches ", 15, INK, False),
       ("Punkt für Punkt", 15, ACCENT_D, True),
       (" mit festen Attributen (Winner, Ace, Doppelfehler, Fehlertypen) und "
        "erzeugen daraus Head-to-Head-Statistiken zur Matchvorbereitung.", 15, INK, False)]],
     line_spacing=1.12)
# four value cards
cards = [
    ("Punkterfassung", "Punkt-für-Punkt live, ITF-konforme Zählung inkl. Tiebreak & Einstand", GREEN_BG, ACCENT_D),
    ("Statistiken", "Head-to-Head, Winner%, Serve%, Aces, Unforced Errors", BLUE_BG, BLUE),
    ("KI-Analyse", "Taktisches Postmortem nach Match-Ende (V1.x)", ORNG_BG, RGBColor(0xE6,0x51,0x00)),
    ("Auth", "Keycloak OAuth2 / OIDC mit PKCE", YELL_BG, RGBColor(0xB7,0x86,0x00)),
]
cw, gap = Inches(2.95), Inches(0.18)
x0 = Inches(0.6)
for i, (ti, de, fill, col) in enumerate(cards):
    x = x0 + i * (cw + gap)
    c = card(s, x, Inches(3.5), cw, Inches(2.3), fill=CARD, shadow=True, line=None)
    tag = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.5), Inches(0.1), Inches(2.3))
    tag.fill.solid(); tag.fill.fore_color.rgb = col; tag.line.fill.background(); tag.shadow.inherit=False
    text(s, x + Inches(0.28), Inches(3.72), cw - Inches(0.45), Inches(0.5),
         [[(ti, 16, INK, True)]])
    text(s, x + Inches(0.28), Inches(4.25), cw - Inches(0.45), Inches(1.4),
         [[(de, 12.5, MUTED, False)]], line_spacing=1.1)
# right info strip
card(s, Inches(7.9), Inches(1.7), Inches(4.8), Inches(1.55), fill=INK, line=None, shadow=True)
text(s, Inches(8.2), Inches(1.92), Inches(4.3), Inches(0.4), [[("TECH-STACK", 11, RGBColor(0x8B,0xE0,0xB0), True)]])
text(s, Inches(8.2), Inches(2.25), Inches(4.3), Inches(1.0),
     [[("Backend  ", 12.5, WHITE, True), ("Spring Boot 4 · Java 25 · Gradle Multi-Module", 12.5, RGBColor(0xC4,0xCD,0xD6), False)],
      [("Frontend  ", 12.5, WHITE, True), ("Angular · Material · ngx-charts", 12.5, RGBColor(0xC4,0xCD,0xD6), False)],
      [("Data/Sec  ", 12.5, WHITE, True), ("PostgreSQL · Keycloak · Docker", 12.5, RGBColor(0xC4,0xCD,0xD6), False)]],
     line_spacing=1.05, space_after=3)

# =============================================================================
# SLIDE 3 — Software-Architektur (Clean Architecture + Module)
# =============================================================================
s = slide()
header(s, "Software-Architektur", "Clean Architecture & Modularer Monolith")
# left: concentric layer diagram
lx, ly = Inches(0.7), Inches(1.85)
layers = [
    (Inches(6.4), Inches(4.6), RGBColor(0xE3,0xEA,0xF0), "Infrastructure  ·  Web (REST), Persistence (JPA), Security"),
    (Inches(5.0), Inches(3.7), RGBColor(0xCF,0xE8,0xDB), "Application  ·  Ports (in/out) + Use-Case-Services"),
    (Inches(3.5), Inches(2.7), ACCENT, "Domain  ·  Model & Rules"),
]
offs = [(0,0),(Inches(0.7),Inches(0.45)),(Inches(1.45),Inches(0.95))]
for (w,h,col,lab),(ox,oy) in zip(layers, offs):
    r = card(s, lx+ox, ly+oy, w, h, fill=col, line=RGBColor(0xC0,0xCB,0xD4), radius=0.04)
for i,(w,h,col,lab) in enumerate(layers):
    ox,oy = offs[i]
    is_core = (i==2)
    tcol = WHITE if is_core else INK
    text(s, lx+ox+Inches(0.2), ly+oy+Inches(0.12), w-Inches(0.4), Inches(0.6),
         [[(lab if not is_core else "Domain", 12.5 if not is_core else 15, tcol, True)]],
         align=PP_ALIGN.CENTER if is_core else PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.MIDDLE if is_core else MSO_ANCHOR.TOP)
text(s, lx, ly+Inches(4.75), Inches(6.4), Inches(0.5),
     [[("Abhängigkeiten zeigen stets nach innen ", 12, MUTED, True),
       ("→ Domain ist framework-frei (Ports & Adapters)", 12, MUTED, False)]],
     align=PP_ALIGN.CENTER)
# right: module list
text(s, Inches(7.55), Inches(1.8), Inches(5.2), Inches(0.4),
     [[("FACHLICHE MODULE (GRADLE MULTI-MODULE)", 11.5, ACCENT_D, True)]])
mods = [
    ("player-module", "Spielerprofile · CRUD & Suche"),
    ("match-module", "Matches, Sets & Scoring (Zählregeln, Punkterfassung)"),
    ("statistics-module", "On-the-fly Stats: H2H, Winner%, Serve%, Aces"),
    ("ai-module", "KI-Match-Analyse über LlmClientPort"),
    ("auth-module", "Keycloak: Token-Validierung & Rollen"),
    ("common-module", "Shared Kernel: DTOs, Exceptions, Config"),
]
my = Inches(2.2)
for i,(name,desc) in enumerate(mods):
    yy = my + i*Inches(0.72)
    card(s, Inches(7.55), yy, Inches(5.18), Inches(0.62), fill=CARD, line=LINE, radius=0.12)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.72), yy+Inches(0.22), Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background(); dot.shadow.inherit=False
    text(s, Inches(8.05), yy+Inches(0.06), Inches(4.6), Inches(0.5),
         [[(name+"   ", 12.5, INK, True), (desc, 11, MUTED, False)]],
         anchor=MSO_ANCHOR.MIDDLE)

# =============================================================================
# SLIDE 4 — System-Architektur (Whitebox / containers + flows)
# =============================================================================
s = slide()
header(s, "System-Architektur", "Whitebox-Gesamtsystem")
# three big blocks: Frontend, Backend(modules), then DB + Keycloak
def block(x,y,w,h,fill,border,title,sub):
    card(s,x,y,w,h,fill=fill,line=border,line_w=1.5,radius=0.05,shadow=True)
    text(s,x+Inches(0.2),y+Inches(0.12),w-Inches(0.4),Inches(0.5),[[(title,14,INK,True)]])
    if sub:
        text(s,x+Inches(0.2),y+Inches(0.5),w-Inches(0.4),Inches(0.5),[[(sub,10.5,MUTED,False)]])

# Frontend
block(Inches(0.6),Inches(2.3),Inches(2.7),Inches(2.1),GREEN_BG,ACCENT_D,
      "Angular SPA","")
text(s,Inches(0.8),Inches(3.0),Inches(2.3),Inches(1.3),
     [[("• touch-optimierte UI",11,INK,False)],[("• REST-Client (Bearer)",11,INK,False)],
      [("• OAuth2 PKCE-Flow",11,INK,False)]],line_spacing=1.1,space_after=3)
# Backend container with module chips
block(Inches(3.75),Inches(1.85),Inches(5.7),Inches(4.5),BLUE_BG,BLUE,
      "Spring Boot Backend (Modularer Monolith)","")
mod_chips = ["player","match + scoring","statistics","ai","auth","common"]
for i,m in enumerate(mod_chips):
    cx = Inches(3.95) + (i%2)*Inches(2.7)
    cy = Inches(2.5) + (i//2)*Inches(0.78)
    chip(s,cx,cy,Inches(2.55),Inches(0.62),[[(m+"-module",12,INK,True)]],
         fill=WHITE,txt_color=INK,line=BLUE)
text(s,Inches(3.95),Inches(5.55),Inches(5.3),Inches(0.7),
     [[("REST /api/*  ·  Use-Case-Services über Ports  ·  synchrone Modulaufrufe via Interfaces",10.5,RGBColor(0x0D,0x47,0x7A),True)]],
     line_spacing=1.05)
# DB + Keycloak
block(Inches(9.95),Inches(1.85),Inches(2.8),Inches(2.0),YELL_BG,RGBColor(0xB7,0x86,0x00),
      "PostgreSQL","")
text(s,Inches(10.15),Inches(2.55),Inches(2.4),Inches(1.2),
     [[("Player · Match · Point",10.5,INK,False)],[("Current_Score · Stats",10.5,INK,False)],
      [("Match_Analysis",10.5,INK,False)]],line_spacing=1.1,space_after=2)
block(Inches(9.95),Inches(4.1),Inches(2.8),Inches(2.0),ORNG_BG,RGBColor(0xE6,0x51,0x00),
      "Keycloak","")
text(s,Inches(10.15),Inches(4.8),Inches(2.4),Inches(1.2),
     [[("OIDC / OAuth2 Issuer",10.5,INK,False)],[("Token-Validierung (JWKS)",10.5,INK,False)],
      [("Google IDP ab V2",10.5,INK,False)]],line_spacing=1.1,space_after=2)
# external LLM
ext = card(s,Inches(3.75),Inches(6.55),Inches(5.7),Inches(0.7),fill=INK,line=None,radius=0.12)
text(s,Inches(3.95),Inches(6.62),Inches(5.3),Inches(0.55),
     [[("LLM-Provider (OpenAI gpt-4o-mini) ",12,WHITE,True),
       ("— extern, über LlmClientPort gekapselt",11,RGBColor(0xC4,0xCD,0xD6),False)]],
     anchor=MSO_ANCHOR.MIDDLE)
# connectors
connector(s,Inches(3.3),Inches(3.35),Inches(3.75),Inches(3.35),BLUE,2.2)
connector(s,Inches(9.45),Inches(2.85),Inches(9.95),Inches(2.85),RGBColor(0xB7,0x86,0x00),2.2)
connector(s,Inches(9.45),Inches(5.1),Inches(9.95),Inches(5.1),RGBColor(0xE6,0x51,0x00),2.2)
connector(s,Inches(6.6),Inches(6.35),Inches(6.6),Inches(6.55),INK,2.0,dash=True)

# =============================================================================
# SLIDE 5 — Deployment / Verteilungssicht
# =============================================================================
s = slide()
header(s, "Deployment", "Verteilungssicht – Docker Compose")
# host frame
card(s, Inches(0.6), Inches(1.75), Inches(8.7), Inches(5.2), fill=WHITE, line=RGBColor(0x9A,0xA6,0xB1), line_w=1.5, radius=0.02)
text(s, Inches(0.8), Inches(1.9), Inches(6), Inches(0.4), [[("🖥  Docker Host  ·  docker-compose", 13, INK, True)]])

def cont(x,y,w,h,fill,border,name,detail,port):
    card(s,x,y,w,h,fill=fill,line=border,line_w=1.25,radius=0.06,shadow=True)
    text(s,x+Inches(0.2),y+Inches(0.14),w-Inches(0.4),Inches(0.45),[[("📦  "+name,13,INK,True)]])
    text(s,x+Inches(0.2),y+Inches(0.58),w-Inches(0.4),Inches(0.6),
         [[(ln,11,MUTED,False)] for ln in detail.split("\n")],line_spacing=1.05,space_after=1)
    chip(s,x+Inches(0.2),y+h-Inches(0.55),Inches(1.7),Inches(0.38),[[(port,10.5,WHITE,True)]],fill=border)

cont(Inches(0.9),Inches(2.5),Inches(3.9),Inches(1.75),GREEN_BG,ACCENT_D,
     "frontend","Nginx + Angular SPA\nproxied /api/ → backend",":80")
cont(Inches(5.05),Inches(2.5),Inches(3.9),Inches(1.75),BLUE_BG,BLUE,
     "backend","Spring Boot API (Java 25)\nnur intern erreichbar",":8080")
cont(Inches(0.9),Inches(4.65),Inches(3.9),Inches(1.9),YELL_BG,RGBColor(0xB7,0x86,0x00),
     "db","PostgreSQL 16\nVolume: postgres-data",":5432")
cont(Inches(5.05),Inches(4.65),Inches(3.9),Inches(1.9),ORNG_BG,RGBColor(0xE6,0x51,0x00),
     "keycloak","Keycloak 26\nHTTPS + HTTP (intern)",":8443 / :18080")
# browser
brow = card(s, Inches(9.6), Inches(2.5), Inches(3.15), Inches(1.0), fill=INK, line=None, radius=0.12, shadow=True)
text(s, Inches(9.6), Inches(2.62), Inches(3.15), Inches(0.8),
     [[("👤  Browser",13,WHITE,True)],[("HTTPS :443 / OAuth2 :8443",10.5,RGBColor(0xC4,0xCD,0xD6),False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2)
# flow notes on the right
text(s, Inches(9.6), Inches(3.85), Inches(3.2), Inches(0.35), [[("VERBINDUNGEN", 11, ACCENT_D, True)]])
bullets(s, Inches(9.6), Inches(4.25), Inches(3.2), Inches(2.6),
        [("Browser → frontend", "  HTTPS :443"),
         ("frontend → backend", "  HTTP-Proxy /api/"),
         ("backend → db", "  JDBC :5432"),
         ("backend → keycloak", "  JWKS :18080"),
         ("Browser → keycloak", "  OAuth2 PKCE :8443")],
        size=11.5, gap=8)
# connectors inside host
connector(s, Inches(4.8), Inches(3.37), Inches(5.05), Inches(3.37), BLUE, 2)
connector(s, Inches(2.85), Inches(4.25), Inches(2.85), Inches(4.65), RGBColor(0xB7,0x86,0x00), 2)
connector(s, Inches(9.6), Inches(3.0), Inches(8.95), Inches(3.2), INK, 1.6, dash=True)

# =============================================================================
# SLIDE 6 — KI-Erkenntnisse & Vorsicht
# =============================================================================
s = slide()
header(s, "KI-Einsatz", "Erkenntnisse & wann ich vorsichtig wäre")
# left column: Erkenntnisse
text(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(0.4),
     [[("✓  WAS GUT FUNKTIONIERT HAT", 12.5, ACCENT_D, True)]])
card(s, Inches(0.6), Inches(2.1), Inches(6.0), Inches(4.85), fill=CARD, line=LINE, shadow=True)
bullets(s, Inches(0.85), Inches(2.35), Inches(5.55), Inches(4.5),
        [("Provider-Abstraktion zahlt sich aus.", " LlmClientPort kapselt das LLM – Wechsel OpenAI→Anthropic/Ollama ohne Eingriff in Use Cases."),
         ("Fake-Adapter statt echtem Call.", " FakeLlmClientAdapter liefert deterministische Ergebnisse: Tests & Entwicklung ohne API-Key und ohne Kosten."),
         ("Strukturierter Output.", " JSON-Schema (BeanOutputConverter) statt fragilem String-Parsing macht die Antwort robust verarbeitbar."),
         ("Kosten & Reproduzierbarkeit by design.", " Analyse 1:1 pro Match persistiert (überschreibbar), Trigger manuell, Mindest-Punktzahl ≥ 10."),
         ("Klare Fehlersemantik.", " 409 / 422 / 502 + persistierter FAILED-Datensatz statt stillem Schlucken von Fehlern.")],
        size=12.5, gap=10)
# right column: Vorsicht
text(s, Inches(6.95), Inches(1.7), Inches(5.8), Inches(0.4),
     [[("⚠  WANN ICH VORSICHTIG WÄRE", 12.5, WARN, True)]])
card(s, Inches(6.95), Inches(2.1), Inches(5.78), Inches(4.85), fill=RGBColor(0xFD,0xF3,0xF3), line=RGBColor(0xF0,0xC6,0xC6), shadow=True)
bullets(s, Inches(7.2), Inches(2.35), Inches(5.3), Inches(4.5),
        [("Nicht-deterministisch.", " Gleicher Input ≠ gleicher Output – ungeeignet, wo exakte/regelbasierte Ergebnisse zählen (z. B. Spielstand-Zählung)."),
         ("Halluzinationen.", " Aussagen wirken plausibel, können aber falsch sein – kritische Ergebnisse müssen verifizierbar/überprüfbar bleiben."),
         ("Datenschutz.", " Keine personenbezogenen Daten ungefiltert an externe LLM senden – DSGVO, Minderjährige im Spielerkontext."),
         ("Kosten & Latenz.", " Jeder Call kostet und dauert; ohne Trigger-Schutz, Caching & Limits skaliert das schlecht."),
         ("Provider-Abhängigkeit.", " API-Änderungen/Ausfälle, Milestone-Stand (Spring AI 2.0-M*) – Fallback & dünner Adapter nötig."),
         ("Mensch entscheidet.", " KI liefert Empfehlungen, nicht Entscheidungen – als Assistenz, nicht als Autorität einsetzen.")],
        size=12, gap=8, marker="•", marker_color=WARN, color=INK)

# =============================================================================
out = "/Users/cbo/Projects/cas/tsas/doc/praesentation/TSaS_Praesentation.pptx"
prs.save(out)
print("saved:", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
